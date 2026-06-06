import fitz
import logging
from docx import Document
from pptx import Presentation

logger = logging.getLogger(__name__)


def extract_docx(path):
    doc = Document(path)
    texts = []
    for p in doc.paragraphs:
        t = p.text.strip()
        if t:
            texts.append(t)
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if row_text:
                texts.append(row_text)
    return texts


def extract_pdf(path):
    doc = fitz.open(path)
    blocks_out = []
    for page in doc:
        blocks = page.get_text("blocks")
        if not blocks:
            continue
        mid_x = page.rect.width / 2
        left_blocks = [b for b in blocks if b[0] < mid_x]
        right_blocks = [b for b in blocks if b[0] >= mid_x]
        if len(left_blocks) > 3 and len(right_blocks) > 3:
            for b in sorted(left_blocks, key=lambda b: (b[1], b[0])):
                t = b[4].strip()
                if t: blocks_out.append(t)
            for b in sorted(right_blocks, key=lambda b: (b[1], b[0])):
                t = b[4].strip()
                if t: blocks_out.append(t)
        else:
            for b in sorted(blocks, key=lambda b: (b[1], b[0])):
                t = b[4].strip()
                if t: blocks_out.append(t)
    return blocks_out


def extract_pptx(path):
    prs = Presentation(path)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        texts.append(t)
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                    if row_text:
                        texts.append(row_text)
    return texts


def extract_text(path):
    try:
        with open(path, encoding="utf-8", errors="replace") as f:
            return f.read()
    except Exception:
        return ""


def extract_xlsx(path):
    from openpyxl import load_workbook
    wb = load_workbook(path, data_only=True)
    rows = []
    for name in wb.sheetnames:
        ws = wb[name]
        rows.append(f"=== {name} ===")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(c.strip() for c in cells):
                rows.append(" | ".join(cells))
    return "\n".join(rows)


def extract_md(path):
    with open(path, encoding="utf-8", errors="replace") as f:
        parts = [p.strip() for p in f.read().split("\n\n") if p.strip()]
        return parts if parts else [open(path, encoding="utf-8", errors="replace").read().strip()]


EXTRACTORS = {
    ".docx": extract_docx,
    ".pdf": extract_pdf,
    ".pptx": extract_pptx,
    ".xlsx": extract_xlsx,
    ".md": extract_md,
}

TEXT_EXTS = {".txt", ".py", ".json", ".yaml", ".yml", ".toml", ".cfg",
             ".ini", ".csv", ".tsv", ".xml", ".html", ".htm", ".tex", ".rst",
             ".log", ".sh", ".bash", ".zsh", ".c", ".h", ".cpp", ".hpp",
             ".java", ".js", ".ts", ".go", ".rs", ".rb", ".lua", ".r",
             ".css", ".sql", ".conf"}


def extract(path, ext):
    try:
        if ext in EXTRACTORS:
            return EXTRACTORS[ext](path)
        if ext in TEXT_EXTS:
            return extract_text(path)
    except Exception:
        logger.exception("parse failed: path=%s ext=%s", path, ext)
        return None
    return None
